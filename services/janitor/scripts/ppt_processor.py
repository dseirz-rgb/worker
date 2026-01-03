#!/usr/bin/env python3
"""
PPT 处理模块

使用 python-pptx 提取 PPT 文件中的文本内容。

使用方法:
    from ppt_processor import process_ppt
    slides = process_ppt("path/to/presentation.pptx")
"""

from dataclasses import dataclass
from typing import List, Optional
from pathlib import Path

from loguru import logger


@dataclass
class SlideContent:
    """幻灯片内容数据结构"""
    page_number: int    # 页码（从 1 开始）
    text: str           # 页面文本（标题 + 正文）
    title: Optional[str] = None  # 幻灯片标题
    
    def to_dict(self) -> dict:
        """转换为字典"""
        return {
            "page_number": self.page_number,
            "text": self.text,
            "title": self.title
        }


def process_ppt(
    file_path: str,
    progress_callback: Optional[callable] = None
) -> List[SlideContent]:
    """
    处理 PPT 文件，返回每页的文本内容
    
    Args:
        file_path: PPT 文件路径 (.pptx)
        progress_callback: 进度回调函数，接收 (current, total, message) 参数
    
    Returns:
        每页幻灯片的内容列表
    
    Raises:
        FileNotFoundError: 文件不存在
        RuntimeError: 处理失败
    """
    # 验证文件存在
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPT 文件不存在: {file_path}")
    
    if path.suffix.lower() != '.pptx':
        logger.warning(f"不支持的文件格式: {path.suffix}，仅支持 .pptx")
    
    logger.info(f"开始处理 PPT: {file_path}")
    
    # 进度回调辅助函数
    def report_progress(current: int, total: int, message: str):
        if progress_callback:
            try:
                progress_callback(current, total, message)
            except Exception as e:
                logger.warning(f"进度回调失败: {e}")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        report_progress(0, 100, "正在打开 PPT 文件...")
        
        # 打开 PPT 文件
        prs = Presentation(file_path)
        
        slides = []
        total_pages = len(prs.slides)
        
        logger.info(f"PPT 共 {total_pages} 页")
        report_progress(10, 100, f"PPT 共 {total_pages} 页，开始提取内容...")
        
        for idx, slide in enumerate(prs.slides, start=1):
            # 提取标题
            title = _extract_title(slide)
            
            # 提取所有文本
            texts = []
            if title:
                texts.append(title)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    # 避免重复添加标题
                    if text and text != title:
                        texts.append(text)
            
            # 合并文本
            page_text = "\n".join(texts)
            
            slides.append(SlideContent(
                page_number=idx,
                text=page_text,
                title=title
            ))
            
            # 报告进度
            progress = 10 + int(80 * idx / total_pages)
            report_progress(progress, 100, f"已处理 {idx}/{total_pages} 页")
            
            logger.debug(f"页 {idx}: {len(page_text)} 字符")
        
        logger.info(f"PPT 处理完成，共提取 {len(slides)} 页内容")
        report_progress(100, 100, f"处理完成，共 {len(slides)} 页")
        return slides
        
    except ImportError:
        logger.error("python-pptx 未安装，请运行: pip install python-pptx")
        raise RuntimeError("python-pptx 未安装")
    except Exception as e:
        logger.error(f"PPT 处理失败: {e}")
        raise RuntimeError(f"PPT 处理失败: {e}")


def _extract_title(slide) -> Optional[str]:
    """
    提取幻灯片标题
    
    Args:
        slide: pptx Slide 对象
    
    Returns:
        标题文本，如果没有则返回 None
    """
    try:
        # 尝试从 title placeholder 获取
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # 尝试从第一个文本框获取
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    return text
        
        return None
    except Exception:
        return None


def get_ppt_metadata(file_path: str) -> dict:
    """
    获取 PPT 文件的元数据
    
    Args:
        file_path: PPT 文件路径
    
    Returns:
        元数据字典
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        return {
            "total_pages": len(prs.slides),
            "file_path": str(Path(file_path).absolute()),
            "file_name": Path(file_path).name,
        }
    except Exception as e:
        logger.error(f"获取 PPT 元数据失败: {e}")
        return {
            "total_pages": 0,
            "file_path": file_path,
            "error": str(e)
        }


def generate_slide_thumbnail(
    file_path: str,
    page_number: int,
    output_path: str,
    width: int = 400
) -> Optional[str]:
    """
    生成 PPT 幻灯片缩略图
    
    注意: 此功能需要 LibreOffice 或 unoconv 支持
    如果不可用，返回 None
    
    Args:
        file_path: PPT 文件路径
        page_number: 页码 (从 1 开始)
        output_path: 输出图片路径
        width: 缩略图宽度
    
    Returns:
        生成的缩略图路径，失败返回 None
    """
    import subprocess
    import tempfile
    
    try:
        # 尝试使用 LibreOffice 转换
        with tempfile.TemporaryDirectory() as tmpdir:
            # 先转换为 PDF
            pdf_path = Path(tmpdir) / "output.pdf"
            
            result = subprocess.run([
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmpdir,
                file_path
            ], capture_output=True, timeout=60)
            
            if result.returncode != 0:
                logger.warning(f"LibreOffice 转换失败: {result.stderr.decode()}")
                return None
            
            # 找到生成的 PDF
            pdf_files = list(Path(tmpdir).glob("*.pdf"))
            if not pdf_files:
                logger.warning("未找到生成的 PDF 文件")
                return None
            
            pdf_path = pdf_files[0]
            
            # 使用 pdftoppm 提取指定页面
            result = subprocess.run([
                "pdftoppm",
                "-png",
                "-f", str(page_number),
                "-l", str(page_number),
                "-scale-to", str(width),
                str(pdf_path),
                str(Path(tmpdir) / "page")
            ], capture_output=True, timeout=30)
            
            if result.returncode != 0:
                logger.warning(f"pdftoppm 转换失败: {result.stderr.decode()}")
                return None
            
            # 找到生成的图片
            png_files = list(Path(tmpdir).glob("page-*.png"))
            if not png_files:
                logger.warning("未找到生成的 PNG 文件")
                return None
            
            # 复制到目标路径
            import shutil
            shutil.copy(png_files[0], output_path)
            
            logger.info(f"缩略图已生成: {output_path}")
            return output_path
            
    except FileNotFoundError:
        logger.warning("LibreOffice 或 pdftoppm 未安装，无法生成缩略图")
        return None
    except subprocess.TimeoutExpired:
        logger.warning("缩略图生成超时")
        return None
    except Exception as e:
        logger.error(f"生成缩略图失败: {e}")
        return None


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("用法: python ppt_processor.py <ppt_file>")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    
    try:
        slides = process_ppt(ppt_file)
        
        print(f"\n处理完成，共 {len(slides)} 页:\n")
        for slide in slides:
            print(f"[页 {slide.page_number}]")
            if slide.title:
                print(f"  标题: {slide.title}")
            print(f"  内容: {slide.text[:100]}{'...' if len(slide.text) > 100 else ''}")
            print()
            
    except Exception as e:
        print(f"错误: {e}")
        sys.exit(1)
